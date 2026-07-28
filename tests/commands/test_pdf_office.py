"""Office-bound test for `mdd pdf` — exercise both Word + PowerPoint legs.

Creates a directory containing one .docx and one .pptx, then invokes
``mdd pdf <dir>`` end-to-end. Only runs on macOS with both Word and
PowerPoint installed.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

import pytest

from mdd.cli import main as cli_main

pytestmark = pytest.mark.office


def _minimal_docx(path: Path) -> None:
    from docx import Document  # pyright: ignore[reportMissingModuleSource]

    doc: Any = Document()  # pyright: ignore[reportAny]
    doc.add_paragraph("Mixed-directory smoke (docx)")
    doc.save(str(path))


def _minimal_pptx(path: Path) -> None:
    import pptx  # pyright: ignore[reportMissingModuleSource]

    prs: Any = pptx.Presentation()  # pyright: ignore[reportAny]
    layout: Any = prs.slide_layouts[0]
    slide: Any = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Mixed-directory smoke (pptx)"
    prs.save(str(path))


def test_pdf_command_handles_mixed_directory(tmp_path: Path) -> None:
    docx = tmp_path / "doc.docx"
    pptx = tmp_path / "deck.pptx"
    _minimal_docx(docx)
    _minimal_pptx(pptx)

    rc = cli_main(["pdf", str(tmp_path)])
    assert rc == 0

    docx_pdf = Path(str(docx) + ".pdf")
    pptx_pdf = Path(str(pptx) + ".pdf")
    assert docx_pdf.exists()
    assert docx_pdf.read_bytes()[:4] == b"%PDF"
    assert pptx_pdf.exists()
    assert pptx_pdf.read_bytes()[:4] == b"%PDF"
