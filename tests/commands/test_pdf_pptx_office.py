"""Office-bound tests for `mdd pdf-pptx` — drive PowerPoint via AppleScript.

These tests build a tiny .pptx with python-pptx and ask the real
``export_pptx_to_pdf`` to drive PowerPoint via AppleScript. They only
run via ``mise run test-office`` (or ``test-all``) on macOS with
Microsoft PowerPoint installed and granted Automation permission for
the running terminal. Marked ``@pytest.mark.office`` so default
``mise run test`` skips them entirely.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

import pytest

from mdd.commands.pdf_pptx import export_pptx_to_pdf, run_pptx_pipeline

pytestmark = pytest.mark.office


def _minimal_pptx(path: Path) -> Path:
    import pptx  # pyright: ignore[reportMissingModuleSource]

    prs: Any = pptx.Presentation()  # pyright: ignore[reportAny]
    layout: Any = prs.slide_layouts[0]
    slide: Any = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Office round-trip smoke"
    prs.save(str(path))
    return path


def test_export_single_pptx_to_pdf(tmp_path: Path) -> None:
    src = _minimal_pptx(tmp_path / "deck.pptx")
    assert export_pptx_to_pdf(src) is True
    pdf = Path(str(src) + ".pdf")
    assert pdf.exists()
    assert pdf.stat().st_size > 0
    # PDF magic bytes
    assert pdf.read_bytes()[:4] == b"%PDF"


def test_pipeline_skips_up_to_date_pdf(tmp_path: Path) -> None:
    src = _minimal_pptx(tmp_path / "deck.pptx")
    # First run: exports the PDF.
    assert run_pptx_pipeline(tmp_path) == 0
    pdf = Path(str(src) + ".pdf")
    mtime1 = pdf.stat().st_mtime
    # Second run: source is not newer, so the pipeline does nothing.
    assert run_pptx_pipeline(tmp_path) == 0
    assert pdf.stat().st_mtime == mtime1
