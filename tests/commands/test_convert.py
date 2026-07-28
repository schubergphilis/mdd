"""Tests for mdd.commands.convert."""

import os
from typing import TYPE_CHECKING, Any
from unittest.mock import MagicMock, patch

import pytest

if TYPE_CHECKING:
    from pathlib import Path

from mdd.cli import main as cli_main
from mdd.commands.convert import (
    SUPPORTED_EXTENSIONS,
    collect_files,
    dest_path,
)
from mdd.converters.docx import (
    convert_body,
    extract_metadata,
    extract_title,
    strip_boilerplate_lines,
    strip_boilerplate_sections,
    strip_empty_headings,
    strip_leading_table,
    strip_toc,
)


def cmd_convert(args: list[str]) -> int:
    """Test helper: invoke `mdd convert` via the argparse entry point."""
    return cli_main(["convert", *args])


def _docx_with_paragraph(tmp_path: Path, name: str, text: str) -> Path:
    from docx import Document  # pyright: ignore[reportMissingModuleSource]

    doc: Any = Document()  # pyright: ignore[reportAny]
    doc.add_paragraph(text)
    p = tmp_path / name
    doc.save(str(p))
    return p


def _docx_with_table(
    tmp_path: Path, name: str, rows: int, cols: int, cells: dict[tuple[int, int], str]
) -> Path:
    from docx import Document  # pyright: ignore[reportMissingModuleSource]

    doc: Any = Document()  # pyright: ignore[reportAny]
    table: Any = doc.add_table(rows=rows, cols=cols)  # pyright: ignore[reportAny]
    for (r, c), val in cells.items():
        table.cell(r, c).text = val
    p = tmp_path / name
    doc.save(str(p))
    return p


class TestStripEmptyHeadings:
    def test_removes_empty_headings(self) -> None:
        lines = ["## ", "### ", "# ", "## Real heading"]
        assert strip_empty_headings(lines) == ["## Real heading"]

    def test_keeps_nonempty_headings(self) -> None:
        lines = ["# Title", "## Section", "content"]
        assert strip_empty_headings(lines) == lines

    def test_strips_artifact_headings(self) -> None:
        lines = ["########### 1 Agenda", "## Real heading"]
        result = strip_empty_headings(lines)
        assert "########### 1 Agenda" not in result
        assert "## Real heading" in result


class TestStripToc:
    def test_removes_word_toc(self) -> None:
        lines = [
            "## Table of Contents",
            "* Introduction 1",
            "* Background 2",
            "## Introduction",
            "Real content here.",
        ]
        result = strip_toc(lines)
        assert "## Table of Contents" not in result
        assert "## Introduction" in result

    def test_removes_docling_tab_toc(self) -> None:
        lines = [
            "########### 1 Agenda",
            "1\tKey points\t2",
            "2\tMeeting Logistics\t2",
            "## 1 Key points",
            "Content here.",
        ]
        result = strip_toc(lines)
        assert "1\tKey points\t2" not in result
        assert "## 1 Key points" in result

    def test_keeps_agenda_section(self) -> None:
        lines = [
            "## Contents",
            "- Deploy the new service",
            "- Review open PRs",
            "## Next section",
        ]
        result = strip_toc(lines)
        assert "## Contents" in result
        assert "- Deploy the new service" in result

    def test_passthrough_when_no_toc(self) -> None:
        lines = ["# Hello", "Some text.", "## Section"]
        assert strip_toc(lines) == lines


class TestStripLeadingTable:
    def test_removes_leading_table(self) -> None:
        md = "| A | B |\n|---|---|\n| x | y |\n\n## Content\nHello."
        assert strip_leading_table(md) == "## Content\nHello."

    def test_no_table_passthrough(self) -> None:
        md = "## Content\nHello."
        assert strip_leading_table(md) == md


class TestStripBoilerplateSections:
    def test_strips_meeting_logistics(self) -> None:
        lines = [
            "## 1 Key points",
            "- Content",
            "## 2 Meeting Logistics",
            "- Boilerplate",
            "## 3 Actions",
        ]
        result = strip_boilerplate_sections(lines)
        assert "## 2 Meeting Logistics" not in result
        assert "## 3 Actions" in result

    def test_passthrough_real_sections(self) -> None:
        lines = ["## 1 Topics", "- item", "## 2 Decisions"]
        assert strip_boilerplate_sections(lines) == lines


class TestStripBoilerplateLines:
    def test_strips_jira_bullet(self) -> None:
        lines = [
            "## Actions",
            "    - Actions are on the [TEC JIRA board](https://jira.example.com) .",
            "- Real item",
        ]
        result = strip_boilerplate_lines(lines)
        assert not any("TEC JIRA board" in ln for ln in result)
        assert "- Real item" in result


class TestDestPath:
    def test_in_place(self, tmp_path: Path) -> None:
        src = tmp_path / "docs" / "report.docx"
        assert dest_path(src, tmp_path, None) == tmp_path / "docs" / "report.docx.md"

    def test_with_dest_dir(self, tmp_path: Path) -> None:
        src_root = tmp_path / "source"
        dest_root = tmp_path / "output"
        src = src_root / "subdir" / "report.docx"
        assert dest_path(src, src_root, dest_root) == dest_root / "subdir" / "report.docx.md"

    def test_dest_dir_with_file_positional_arg(self, tmp_path: Path) -> None:
        """#2: --dest-dir combined with a file positional arg must not produce '.md'."""
        src = _docx_with_paragraph(tmp_path, "simple-document.docx", "Hello")
        dest_dir = tmp_path / "out"
        dest_dir.mkdir()
        with patch("mdd.converters.docx.convert_body", return_value="# Hello"):
            result = cmd_convert(["--dest-dir", str(dest_dir), str(src)])
        assert result == 0
        output_file = dest_dir / "simple-document.docx.md"
        assert output_file.exists(), f"Expected {output_file}, got: {list(dest_dir.iterdir())}"
        assert not (dest_dir / ".md").exists(), "Bug: produced empty-stem '.md' file"


class TestExtractTitle:
    def test_returns_empty_for_plain_docx(self, tmp_path: Path) -> None:
        p = _docx_with_paragraph(tmp_path, "plain.docx", "Hello")
        assert extract_title(p) == ""

    def test_returns_core_properties_title(self, tmp_path: Path) -> None:
        from docx import Document  # pyright: ignore[reportMissingModuleSource]

        doc: Any = Document()  # pyright: ignore[reportAny]
        doc.core_properties.title = "My Report"
        p = tmp_path / "titled.docx"
        doc.save(str(p))
        assert extract_title(p) == "My Report"


class TestExtractMetadata:
    def test_returns_empty_when_no_table(self, tmp_path: Path) -> None:
        p = _docx_with_paragraph(tmp_path, "test.docx", "Hello world")
        assert extract_metadata(p) == {}

    def test_extracts_from_two_column_table(self, tmp_path: Path) -> None:
        p = _docx_with_table(
            tmp_path,
            "meta.docx",
            rows=3,
            cols=2,
            cells={
                (0, 0): "Date",
                (0, 1): "2025-10-13",
                (1, 0): "Location",
                (1, 1): "Hybrid",
                (2, 0): "Participants",
                (2, 1): "Alice, Bob",
            },
        )
        result = extract_metadata(p)
        assert result["Date"] == "2025-10-13"
        assert result["Location"] == "Hybrid"

    def test_normalises_invited_to_invitees(self, tmp_path: Path) -> None:
        p = _docx_with_table(
            tmp_path, "inv.docx", rows=1, cols=2, cells={(0, 0): "Invited", (0, 1): "Carol"}
        )
        assert "Invitees" in extract_metadata(p)

    def test_skips_tables_over_20_rows(self, tmp_path: Path) -> None:
        p = _docx_with_table(
            tmp_path, "big.docx", rows=21, cols=2, cells={(0, 0): "Date", (0, 1): "2025-01-01"}
        )
        assert extract_metadata(p) == {}


class TestCmdConvert:
    def test_no_args_prints_help(self, capsys: pytest.CaptureFixture[str]) -> None:
        assert cmd_convert([]) == 1

    def test_nonexistent_path_exits_nonzero(
        self, tmp_path: Path, capsys: pytest.CaptureFixture[str]
    ) -> None:
        result = cmd_convert([str(tmp_path / "does-not-exist")])
        assert result == 1
        assert "does not exist" in capsys.readouterr().err

    def test_nonexistent_file_flag_exits_nonzero(
        self, tmp_path: Path, capsys: pytest.CaptureFixture[str]
    ) -> None:
        result = cmd_convert(["--file", str(tmp_path / "no-such-file.docx")])
        assert result == 1
        assert "does not exist" in capsys.readouterr().err

    def test_help_flag(self, capsys: pytest.CaptureFixture[str]) -> None:
        with pytest.raises(SystemExit) as exc_info:
            cmd_convert(["--help"])
        assert exc_info.value.code == 0

    def test_dry_run_does_not_write(self, tmp_path: Path) -> None:
        _docx_with_paragraph(tmp_path, "test.docx", "Hello")
        with patch("mdd.converters.docx.convert_body", return_value="# Hello"):
            cmd_convert(["--dry-run", str(tmp_path)])
        assert not (tmp_path / "test.docx.md").exists()

    def test_skips_up_to_date_files(self, tmp_path: Path) -> None:
        docx_path = _docx_with_paragraph(tmp_path, "test.docx", "Content")
        md_path = tmp_path / "test.docx.md"
        md_path.write_text("already converted")
        mtime = docx_path.stat().st_mtime
        os.utime(md_path, (mtime + 1, mtime + 1))

        mock = MagicMock(return_value="# Body")
        with patch("mdd.converters.docx.convert_body", mock):
            cmd_convert([str(tmp_path)])
        mock.assert_not_called()

    def test_force_regenerates(self, tmp_path: Path) -> None:
        docx_path = _docx_with_paragraph(tmp_path, "test.docx", "Content")
        md_path = tmp_path / "test.docx.md"
        md_path.write_text("stale")
        mtime = docx_path.stat().st_mtime
        os.utime(md_path, (mtime + 1, mtime + 1))

        mock = MagicMock(return_value="# Regenerated")
        with patch("mdd.converters.docx.convert_body", mock):
            cmd_convert(["--force", str(tmp_path)])
        mock.assert_called_once()

    def test_failure_exits_nonzero(
        self, tmp_path: Path, capsys: pytest.CaptureFixture[str]
    ) -> None:
        _docx_with_paragraph(tmp_path, "test.docx", "Content")
        with patch("mdd.converters.docx.convert_body", side_effect=RuntimeError("boom")):
            result = cmd_convert([str(tmp_path)])
        assert result == 1
        assert "boom" in capsys.readouterr().err


class TestConvertEndToEnd:
    def test_convert_single_docx(self, tmp_path: Path) -> None:
        fixture = _docx_with_paragraph(tmp_path, "sample.docx", "Hello world")
        result = convert_body(fixture)
        assert result.strip()
        assert "<" not in result
        assert ">" not in result


def _make_pptx_file(tmp_path: Path, name: str = "test.pptx") -> Path:
    """Create a minimal .pptx file using python-pptx."""
    import pptx  # pyright: ignore[reportMissingModuleSource]

    prs: Any = pptx.Presentation()  # pyright: ignore[reportAny]
    slide_layout: Any = prs.slide_layouts[0]
    slide: Any = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Test Slide"
    path = tmp_path / name
    prs.save(str(path))
    return path


class TestCollectFiles:
    def test_collects_docx(self, tmp_path: Path) -> None:
        p = _docx_with_paragraph(tmp_path, "doc.docx", "text")
        result = collect_files(tmp_path)
        assert p in result

    def test_collects_pptx(self, tmp_path: Path) -> None:
        p = _make_pptx_file(tmp_path, "deck.pptx")
        result = collect_files(tmp_path)
        assert p in result

    def test_collects_single_file(self, tmp_path: Path) -> None:
        p = _make_pptx_file(tmp_path, "single.pptx")
        result = collect_files(p)
        assert result == [p]

    def test_ignores_unsupported_extensions(self, tmp_path: Path) -> None:
        (tmp_path / "ignore.xlsx").write_text("spreadsheet")
        result = collect_files(tmp_path)
        assert all(f.suffix.lower() in SUPPORTED_EXTENSIONS for f in result)

    def test_supported_extensions_contains_expected(self) -> None:
        assert ".docx" in SUPPORTED_EXTENSIONS
        assert ".pptx" in SUPPORTED_EXTENSIONS
        assert ".pdf" in SUPPORTED_EXTENSIONS

    def test_doc_not_in_supported_extensions(self) -> None:
        """#48: .doc is not supported; removed from SUPPORTED_EXTENSIONS."""
        assert ".doc" not in SUPPORTED_EXTENSIONS


class TestDocLegacy:
    """#48: .doc files should be rejected with a helpful message."""

    def test_doc_file_exits_nonzero(
        self, tmp_path: Path, capsys: pytest.CaptureFixture[str]
    ) -> None:
        doc_path = tmp_path / "legacy.doc"
        doc_path.write_bytes(b"\xd0\xcf\x11\xe0")  # OLE compound doc magic bytes
        result = cmd_convert(["--file", str(doc_path)])
        assert result == 1
        err = capsys.readouterr().err
        assert "legacy .doc" in err.lower() or ".doc" in err

    def test_doc_file_not_collected(self, tmp_path: Path) -> None:
        doc_path = tmp_path / "old.doc"
        doc_path.write_bytes(b"\xd0\xcf\x11\xe0")
        result = collect_files(tmp_path)
        assert doc_path not in result


class TestCmdConvertPptx:
    def test_pptx_dry_run_does_not_write(self, tmp_path: Path) -> None:
        p = _make_pptx_file(tmp_path, "deck.pptx")
        cmd_convert(["--dry-run", str(p)])
        assert not (tmp_path / "deck.pptx.md").exists()

    def test_pptx_conversion_via_cmd(self, tmp_path: Path) -> None:
        p = _make_pptx_file(tmp_path, "deck.pptx")
        mock = MagicMock()
        with patch("mdd.convert.pptx.convert_pptx", mock):
            result = cmd_convert([str(p)])
        assert result == 0
        mock.assert_called_once()
        call_src = mock.call_args[0][0]
        assert call_src == p

    def test_pptx_skips_up_to_date(self, tmp_path: Path) -> None:
        p = _make_pptx_file(tmp_path, "deck.pptx")
        md_path = tmp_path / "deck.pptx.md"
        md_path.write_text("already done")
        mtime = p.stat().st_mtime
        os.utime(md_path, (mtime + 1, mtime + 1))

        mock = MagicMock()
        with patch("mdd.convert.pptx.convert_pptx", mock):
            cmd_convert([str(tmp_path)])
        mock.assert_not_called()

    def test_pptx_force_regenerates(self, tmp_path: Path) -> None:
        p = _make_pptx_file(tmp_path, "deck.pptx")
        md_path = tmp_path / "deck.pptx.md"
        md_path.write_text("stale")
        mtime = p.stat().st_mtime
        os.utime(md_path, (mtime + 1, mtime + 1))

        mock = MagicMock()
        with patch("mdd.convert.pptx.convert_pptx", mock):
            cmd_convert(["--force", str(tmp_path)])
        mock.assert_called_once()

    def test_pptx_via_file_flag(self, tmp_path: Path) -> None:
        p = _make_pptx_file(tmp_path, "single.pptx")
        mock = MagicMock()
        with patch("mdd.convert.pptx.convert_pptx", mock):
            result = cmd_convert(["--file", str(p)])
        assert result == 0
        mock.assert_called_once()
