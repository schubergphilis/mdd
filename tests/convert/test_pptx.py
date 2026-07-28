"""Tests for mdd.convert.pptx — convert_pptx function."""

import io
from typing import TYPE_CHECKING, Any

import pytest

if TYPE_CHECKING:
    from collections.abc import Callable
    from pathlib import Path


def _make_pptx(  # noqa: C901, PLR0912, PLR0915
    tmp_path: Path,
    slides: list[dict[str, Any]],
    *,
    filename: str = "test.pptx",
    core_title: str = "",
    core_author: str = "",
) -> Path:
    """Build a .pptx in tmp_path from a list of slide descriptors.

    Each slide dict may have:
      title (str): slide title text
      bullets (list[tuple[int,str]]): (level, text) pairs
      notes (str): speaker notes
      pictures (list[bytes]): raw PNG bytes to embed as pictures
      table (list[list[str]]): rows x cols of text for a table
    """
    import pptx  # pyright: ignore[reportMissingModuleSource]
    from pptx.util import Inches  # pyright: ignore[reportMissingModuleSource]

    prs: Any = pptx.Presentation()  # pyright: ignore[reportAny]
    if core_title:
        prs.core_properties.title = core_title
    if core_author:
        prs.core_properties.author = core_author

    blank_layout: Any = prs.slide_layouts[6]  # blank layout
    title_content_layout: Any = prs.slide_layouts[1]  # title + content

    for slide_def in slides:
        slide_title: str = slide_def.get("title", "")
        bullets: list[tuple[int, str]] = slide_def.get("bullets", [])
        notes_text: str = slide_def.get("notes", "")
        pictures: list[bytes] = slide_def.get("pictures", [])
        table_data: list[list[str]] = slide_def.get("table", [])

        if slide_title or bullets:
            slide: Any = prs.slides.add_slide(title_content_layout)
            if slide_title:
                slide.shapes.title.text = slide_title
            if bullets:
                # Use the content placeholder (index 1)
                content_ph: Any = slide.placeholders[1]
                tf: Any = content_ph.text_frame
                tf.clear()
                for i, (level, text) in enumerate(bullets):
                    if i == 0:
                        para: Any = tf.paragraphs[0]
                    else:
                        para = tf.add_paragraph()
                    para.text = text
                    para.level = level
        else:
            slide = prs.slides.add_slide(blank_layout)

        # Add pictures
        for pic_bytes in pictures:
            img_stream = io.BytesIO(pic_bytes)
            slide.shapes.add_picture(img_stream, Inches(1), Inches(1), Inches(2), Inches(2))

        # Add table
        if table_data and len(table_data) > 0:
            rows = len(table_data)
            cols = len(table_data[0]) if table_data else 1
            tbl_shape: Any = slide.shapes.add_table(  # pyright: ignore[reportAny]
                rows, cols, Inches(1), Inches(4), Inches(6), Inches(2)
            )
            tbl: Any = tbl_shape.table
            for r, row in enumerate(table_data):
                for c, cell_text in enumerate(row):
                    tbl.cell(r, c).text = cell_text

        # Add speaker notes
        if notes_text:
            notes_slide: Any = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text

    path = tmp_path / filename
    prs.save(str(path))
    return path


def _minimal_png() -> bytes:
    """Return minimal valid 1x1 PNG bytes."""
    # Hand-crafted 1x1 transparent PNG
    import base64

    png_b64 = (
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk"
        "+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
    )
    return base64.b64decode(png_b64)


class TestConvertPptxTitleOnly:
    def test_produces_h2_heading(self, tmp_path: Path) -> None:
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(tmp_path, [{"title": "My Slide"}])
        dst = tmp_path / "test.pptx.md"
        convert_pptx(src, dst)
        content = dst.read_text()
        assert "## My Slide" in content

    def test_no_extra_body_lines(self, tmp_path: Path) -> None:
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(tmp_path, [{"title": "Only Title"}])
        dst = tmp_path / "test.pptx.md"
        convert_pptx(src, dst)
        content = dst.read_text()
        # Should have frontmatter and heading, but no bullet content (- item lines)
        # Exclude YAML fence lines (---) which also start with -
        bullet_lines = [
            line
            for line in content.splitlines()
            if line.strip().startswith("-") and not line.strip().startswith("---")
        ]
        assert bullet_lines == []

    def test_title_not_duplicated_as_body(self, tmp_path: Path) -> None:
        """Regression: title placeholder must not also be emitted as body text.

        python-pptx returns fresh wrapper objects on each shape lookup, so
        identity (``is``) comparison between ``slide.shapes.title`` and the
        shape yielded by iterating ``slide.shapes`` is False even when both
        wrap the same XML element. The old code used ``is`` and consequently
        emitted the title text twice — once as ``## Title`` and once as a
        plain-text paragraph in the body.
        """
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(tmp_path, [{"title": "Example PowerPoint"}])
        dst = tmp_path / "test.pptx.md"
        convert_pptx(src, dst)
        content = dst.read_text()
        # The title text should appear exactly once — as the ## heading.
        assert content.count("Example PowerPoint") == 1
        assert "## Example PowerPoint" in content


class TestConvertPptxBullets:
    def test_level0_bullets_as_plain(self, tmp_path: Path) -> None:
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(
            tmp_path,
            [{"title": "Bullets", "bullets": [(0, "Top level"), (1, "Sub item")]}],
        )
        dst = tmp_path / "test.pptx.md"
        convert_pptx(src, dst)
        content = dst.read_text()
        assert "Top level" in content
        assert "- Sub item" in content

    def test_multiple_levels(self, tmp_path: Path) -> None:
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(
            tmp_path,
            [
                {
                    "title": "Multi",
                    "bullets": [
                        (0, "Level zero"),
                        (1, "Level one"),
                        (2, "Level two"),
                    ],
                }
            ],
        )
        dst = tmp_path / "test.pptx.md"
        convert_pptx(src, dst)
        content = dst.read_text()
        assert "Level zero" in content
        assert "- Level one" in content
        assert "  - Level two" in content


class TestConvertPptxNotes:
    def test_notes_block_appears(self, tmp_path: Path) -> None:
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(tmp_path, [{"title": "Noted", "notes": "Speaker note here"}])
        dst = tmp_path / "test.pptx.md"
        convert_pptx(src, dst)
        content = dst.read_text()
        assert "::: notes" in content
        assert "Speaker note here" in content
        assert ":::" in content

    def test_no_notes_block_when_empty(self, tmp_path: Path) -> None:
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(tmp_path, [{"title": "No notes"}])
        dst = tmp_path / "test.pptx.md"
        convert_pptx(src, dst)
        content = dst.read_text()
        assert "::: notes" not in content


class TestConvertPptxPictures:
    def test_picture_written_to_attachments(self, tmp_path: Path) -> None:
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(tmp_path, [{"title": "Pics", "pictures": [_minimal_png()]}])
        dst = tmp_path / "test.pptx.md"
        convert_pptx(src, dst)
        attachments = tmp_path / "test.pptx-attachments"
        assert attachments.is_dir()
        img_files = list(attachments.iterdir())
        assert len(img_files) == 1

    def test_picture_link_in_markdown(self, tmp_path: Path) -> None:
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(tmp_path, [{"title": "Pics", "pictures": [_minimal_png()]}])
        dst = tmp_path / "test.pptx.md"
        convert_pptx(src, dst)
        content = dst.read_text()
        # Filenames are now content-addressed (sha1(blob)[:16]) — see #93.
        assert "![](test.pptx-attachments/image_" in content

    def test_same_blob_dedup_across_slides(self, tmp_path: Path) -> None:
        """#93: identical blobs across slides write one file and share the link."""
        from mdd.convert.pptx import convert_pptx

        png = _minimal_png()
        src = _make_pptx(
            tmp_path,
            [
                {"title": "Slide 1", "pictures": [png]},
                {"title": "Slide 2", "pictures": [png]},
            ],
        )
        dst = tmp_path / "test.pptx.md"
        convert_pptx(src, dst)
        attachments = tmp_path / "test.pptx-attachments"
        files = list(attachments.iterdir())
        # Same blob -> one file, not two.
        assert len(files) == 1
        # And both slides reference that one file.
        content = dst.read_text()
        link_count = content.count(f"](test.pptx-attachments/{files[0].name})")
        assert link_count == 2

    def test_png_extension_preserved(self, tmp_path: Path) -> None:
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(tmp_path, [{"title": "Ext", "pictures": [_minimal_png()]}])
        dst = tmp_path / "test.pptx.md"
        convert_pptx(src, dst)
        attachments = tmp_path / "test.pptx-attachments"
        names = [f.name for f in attachments.iterdir()]
        assert any(n.endswith(".png") for n in names)


class TestConvertPptxUntitled:
    def test_fallback_slide_n_heading(self, tmp_path: Path) -> None:
        from mdd.convert.pptx import convert_pptx

        # No title in slide_def -> uses blank layout -> no title shape
        src = _make_pptx(tmp_path, [{}])
        dst = tmp_path / "test.pptx.md"
        convert_pptx(src, dst)
        content = dst.read_text()
        assert "## Slide 1" in content

    def test_multiple_untitled_slides_numbered(self, tmp_path: Path) -> None:
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(tmp_path, [{}, {}])
        dst = tmp_path / "test.pptx.md"
        convert_pptx(src, dst)
        content = dst.read_text()
        assert "## Slide 1" in content
        assert "## Slide 2" in content


class TestConvertPptxTable:
    def test_simple_table_as_markdown(self, tmp_path: Path) -> None:
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(
            tmp_path,
            [
                {
                    "table": [
                        ["Name", "Value"],
                        ["Alice", "42"],
                        ["Bob", "7"],
                    ]
                }
            ],
        )
        dst = tmp_path / "test.pptx.md"
        convert_pptx(src, dst)
        content = dst.read_text()
        assert "| Name" in content
        assert "| Alice" in content
        assert "---" in content

    def test_table_header_separator_present(self, tmp_path: Path) -> None:
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(
            tmp_path,
            [{"table": [["Col A", "Col B"], ["x", "y"]]}],
        )
        dst = tmp_path / "test.pptx.md"
        convert_pptx(src, dst)
        content = dst.read_text()
        lines = content.splitlines()
        # Find the separator line (all dashes and pipes)
        sep_lines = [
            line
            for line in lines
            if set(line.replace("|", "").replace("-", "").replace(" ", "")) == set()
        ]
        assert sep_lines, "Expected a table separator row"


class TestConvertPptxFrontmatter:
    def test_frontmatter_has_slide_count(self, tmp_path: Path) -> None:
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(
            tmp_path,
            [{"title": "A"}, {"title": "B"}, {"title": "C"}],
        )
        dst = tmp_path / "test.pptx.md"
        convert_pptx(src, dst)
        content = dst.read_text()
        assert "slide_count: 3" in content

    def test_frontmatter_has_title_when_set(self, tmp_path: Path) -> None:
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(
            tmp_path,
            [{"title": "Slide"}],
            core_title="Deck Title",
        )
        dst = tmp_path / "test.pptx.md"
        convert_pptx(src, dst)
        content = dst.read_text()
        assert "Deck Title" in content

    def test_frontmatter_fences(self, tmp_path: Path) -> None:
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(tmp_path, [{"title": "X"}])
        dst = tmp_path / "test.pptx.md"
        convert_pptx(src, dst)
        content = dst.read_text()
        assert content.startswith("---")
        assert "pptx:" in content


class TestRenderTableHtmlEscaping:
    """#23: HTML table fallback must escape cell text."""

    def test_script_tag_escaped_in_html_table(self, tmp_path: Path) -> None:
        from mdd.convert.pptx import _render_table_html  # pyright: ignore[reportPrivateUsage]

        # Build a mock table with a cell containing XSS payload
        class FakeCell:
            class FakeFrame:
                text = "<script>alert(1)</script>"

            text_frame = FakeFrame()

        class FakeRow:
            cells = [FakeCell(), FakeCell()]  # noqa: RUF012

        class FakeTable:
            rows = [FakeRow(), FakeRow()]  # noqa: RUF012

        result = _render_table_html(FakeTable())
        html_block = "\n".join(result)
        assert "<script>" not in html_block
        assert "&lt;script&gt;" in html_block

    def test_html_entity_chars_escaped(self, tmp_path: Path) -> None:
        from mdd.convert.pptx import _render_table_html  # pyright: ignore[reportPrivateUsage]

        class FakeCell:
            class FakeFrame:
                text = "A & B > C"

            text_frame = FakeFrame()

        class FakeRow:
            cells = [FakeCell()]  # noqa: RUF012

        class FakeTable:
            rows = [FakeRow(), FakeRow()]  # noqa: RUF012

        result = _render_table_html(FakeTable())
        html_block = "\n".join(result)
        assert "A & B" not in html_block
        assert "&amp;" in html_block


class TestNotesTextExceptions:
    """#69: _notes_text should only catch AttributeError/KeyError, not all exceptions."""

    def test_propagates_value_error(self) -> None:
        from mdd.convert.pptx import _notes_text  # pyright: ignore[reportPrivateUsage]

        class BadNotesSlide:
            @property
            def notes_text_frame(self) -> None:
                raise ValueError("unexpected xml")

        class FakeSlide:
            @property
            def notes_slide(self) -> BadNotesSlide:
                return BadNotesSlide()

        with pytest.raises(ValueError, match="unexpected xml"):
            _notes_text(FakeSlide())

    def test_attribute_error_returns_empty(self) -> None:
        from mdd.convert.pptx import _notes_text  # pyright: ignore[reportPrivateUsage]

        class FakeSlide:
            @property
            def notes_slide(self) -> None:
                raise AttributeError("no notes_slide")

        assert _notes_text(FakeSlide()) == ""

    def test_key_error_returns_empty(self) -> None:
        from mdd.convert.pptx import _notes_text  # pyright: ignore[reportPrivateUsage]

        class BadNotesSlide:
            @property
            def notes_text_frame(self) -> None:
                raise KeyError("missing element")

        class FakeSlide:
            @property
            def notes_slide(self) -> BadNotesSlide:
                return BadNotesSlide()

        assert _notes_text(FakeSlide()) == ""


class TestFrontmatterYamlEncoding:
    """#24: Frontmatter must use proper YAML encoding, not repr()."""

    def test_title_with_apostrophe(self, tmp_path: Path) -> None:
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(tmp_path, [{"title": "Slide"}], core_title="It's Fine")
        dst = tmp_path / "test.pptx.md"
        convert_pptx(src, dst)
        content = dst.read_text()
        # repr() would produce `'It\'s Fine'` which is invalid YAML
        # yaml.safe_dump produces `"It's Fine"` (double-quoted)
        import yaml

        lines = content.splitlines()
        fm_end = lines.index("---", 1)
        fm_text = "\n".join(lines[1:fm_end])
        data = yaml.safe_load(fm_text)
        assert data["pptx"]["title"] == "It's Fine"

    def test_title_with_double_quote(self, tmp_path: Path) -> None:
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(tmp_path, [{"title": "X"}], core_title='Say "hello"')
        dst = tmp_path / "test.pptx.md"
        convert_pptx(src, dst)
        content = dst.read_text()
        import yaml

        lines = content.splitlines()
        fm_end = lines.index("---", 1)
        fm_text = "\n".join(lines[1:fm_end])
        data = yaml.safe_load(fm_text)
        assert data["pptx"]["title"] == 'Say "hello"'


class TestConvertPptxAtomicWrite:
    def test_no_tmp_file_after_conversion(self, tmp_path: Path) -> None:
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(tmp_path, [{"title": "Atomic"}])
        dst = tmp_path / "test.pptx.md"
        convert_pptx(src, dst)
        tmp_file = dst.with_suffix(".md.tmp")
        assert not tmp_file.exists()
        assert dst.exists()

    def test_converts_real_pptx_roundtrip(self, tmp_path: Path) -> None:
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(
            tmp_path,
            [
                {"title": "Intro", "bullets": [(0, "Point one"), (1, "Sub point")]},
                {"title": "Details", "notes": "These are speaker notes"},
            ],
            core_title="Full Deck",
            core_author="Test Author",
        )
        dst = tmp_path / "test.pptx.md"
        convert_pptx(src, dst)
        content = dst.read_text()
        assert "## Intro" in content
        assert "Point one" in content
        assert "## Details" in content
        assert "::: notes" in content
        assert "speaker notes" in content


class TestUnrecognizedImageDrop:
    """#91: unsupported image formats (e.g. MPO) drop, never abort the deck."""

    @staticmethod
    def _patch_image_ext_raises(
        monkeypatch: pytest.MonkeyPatch,
        exc_factory: Callable[[int], BaseException | None],
    ) -> dict[str, int]:
        """Force ``Image.ext`` to raise; return a call-count dict.

        ``exc_factory(call_count)`` is invoked on every access of
        ``Image.ext``; returning ``None`` falls through to the original
        ``lazyproperty`` lookup so callers can opt-in to "first call
        raises, rest are real".
        """
        from pptx.parts.image import Image  # pyright: ignore[reportMissingModuleSource]

        counter = {"n": 0}
        # The class attribute is a `lazyproperty` descriptor at runtime;
        # the type stubs declare `ext: str`, hence the Any cast.
        original_descriptor: Any = Image.__dict__["ext"]  # pyright: ignore[reportAny,reportExplicitAny]

        def _ext_getter(self: Any) -> Any:  # pyright: ignore[reportAny,reportExplicitAny]
            counter["n"] += 1
            exc = exc_factory(counter["n"])
            if exc is not None:
                raise exc
            return original_descriptor.__get__(self, type(self))  # pyright: ignore[reportAny]

        monkeypatch.setattr(Image, "ext", property(_ext_getter))
        return counter

    @staticmethod
    def _mpo_value_error() -> ValueError:
        return ValueError(
            "unsupported image format, expected one of: "
            "dict_keys(['BMP', 'GIF', 'JPEG', 'PNG', 'TIFF', 'WMF']), got 'MPO'"
        )

    def test_unrecognized_format_does_not_raise(
        self,
        tmp_path: Path,
        monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        from mdd.convert.pptx import convert_pptx

        # Two slides with a normal PNG; we make python-pptx's Image.ext
        # raise on the first access (MPO/HEIC-shaped failure) and succeed
        # thereafter.
        src = _make_pptx(
            tmp_path,
            [
                {"title": "MPO slide", "pictures": [_minimal_png()]},
                {"title": "OK slide", "pictures": [_minimal_png()]},
            ],
        )
        dst = tmp_path / "test.pptx.md"

        def _factory(n: int) -> BaseException | None:
            return self._mpo_value_error() if n == 1 else None

        self._patch_image_ext_raises(monkeypatch, _factory)

        # Must not raise
        convert_pptx(src, dst)
        assert dst.exists()

    def test_dropped_image_placeholder_in_markdown(
        self,
        tmp_path: Path,
        monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(tmp_path, [{"title": "MPO", "pictures": [_minimal_png()]}])
        dst = tmp_path / "test.pptx.md"

        def _factory(_n: int) -> BaseException:
            return self._mpo_value_error()

        self._patch_image_ext_raises(monkeypatch, _factory)
        convert_pptx(src, dst)
        content = dst.read_text()
        assert "dropped unsupported image: MPO" in content

    def test_drop_summary_emitted(
        self,
        tmp_path: Path,
        monkeypatch: pytest.MonkeyPatch,
        caplog: pytest.LogCaptureFixture,
    ) -> None:
        import logging

        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(tmp_path, [{"title": "MPO", "pictures": [_minimal_png()]}])
        dst = tmp_path / "test.pptx.md"

        def _factory(_n: int) -> BaseException:
            return self._mpo_value_error()

        self._patch_image_ext_raises(monkeypatch, _factory)
        mdd_log = logging.getLogger("mdd.convert.pptx")
        mdd_log.addHandler(caplog.handler)
        try:
            with caplog.at_level(logging.WARNING, logger="mdd.convert.pptx"):
                convert_pptx(src, dst)
        finally:
            mdd_log.removeHandler(caplog.handler)
        msgs = " ".join(r.getMessage() for r in caplog.records)
        assert "dropped 1 embedded image" in msgs
        assert "MPO" in msgs

    def test_io_errors_still_propagate(
        self,
        tmp_path: Path,
        monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        """A non-format I/O error must not be silently swallowed."""
        from mdd.convert import pptx as pptx_mod
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(tmp_path, [{"title": "Boom", "pictures": [_minimal_png()]}])
        dst = tmp_path / "test.pptx.md"

        # Simulate a disk-side failure inside write_image (which is the
        # post-`image.ext` seam where genuine I/O errors can surface).
        def fake_write_image(*_args: Any, **_kwargs: Any) -> Any:
            raise OSError("disk on fire")

        monkeypatch.setattr("mdd.convert.images.write_image", fake_write_image)
        # Reference the module under test to keep the import live.
        assert pptx_mod is not None
        with pytest.raises(OSError, match="disk on fire"):
            convert_pptx(src, dst)


class TestShapeImageAccessFailures:
    """#128: ``shape.image`` itself can raise; widen catch so the deck still converts.

    Two upstream python-pptx failure modes propagate past the (now-wider)
    try/except in ``_extract_picture``:

    - ``ValueError("no embedded image")`` — picture-shape image part missing.
    - ``AttributeError("'Part' object has no attribute 'image'")`` — picture-shape
      relationship target is a non-image Part (e.g. embedded OLE object).

    Both must be classified as ``unreadable`` and recorded in
    ``dropped_image_reasons`` rather than aborting the conversion.
    """

    @staticmethod
    def _patch_picture_image_raises(
        monkeypatch: pytest.MonkeyPatch,
        exc: BaseException,
    ) -> None:
        """Force ``Picture.image`` to raise ``exc`` on every access."""
        from pptx.shapes.picture import Picture  # pyright: ignore[reportMissingModuleSource]

        def _image_getter(_self: Any) -> Any:  # pyright: ignore[reportAny,reportExplicitAny]
            raise exc

        monkeypatch.setattr(Picture, "image", property(_image_getter))

    def test_no_embedded_image_value_error_does_not_raise(
        self,
        tmp_path: Path,
        monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(tmp_path, [{"title": "Bad pic", "pictures": [_minimal_png()]}])
        dst = tmp_path / "test.pptx.md"
        self._patch_picture_image_raises(monkeypatch, ValueError("no embedded image"))

        # Must not raise — the whole deck used to abort here.
        convert_pptx(src, dst)
        assert dst.exists()

    def test_no_embedded_image_emits_unreadable_placeholder(
        self,
        tmp_path: Path,
        monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(tmp_path, [{"title": "Bad pic", "pictures": [_minimal_png()]}])
        dst = tmp_path / "test.pptx.md"
        self._patch_picture_image_raises(monkeypatch, ValueError("no embedded image"))

        convert_pptx(src, dst)
        content = dst.read_text()
        assert "dropped unsupported image: unreadable" in content

    def test_part_has_no_image_attribute_error_does_not_raise(
        self,
        tmp_path: Path,
        monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(tmp_path, [{"title": "OLE pic", "pictures": [_minimal_png()]}])
        dst = tmp_path / "test.pptx.md"
        self._patch_picture_image_raises(
            monkeypatch,
            AttributeError("'Part' object has no attribute 'image'"),
        )

        convert_pptx(src, dst)
        assert dst.exists()

    def test_part_has_no_image_emits_unreadable_placeholder(
        self,
        tmp_path: Path,
        monkeypatch: pytest.MonkeyPatch,
    ) -> None:
        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(tmp_path, [{"title": "OLE pic", "pictures": [_minimal_png()]}])
        dst = tmp_path / "test.pptx.md"
        self._patch_picture_image_raises(
            monkeypatch,
            AttributeError("'Part' object has no attribute 'image'"),
        )

        convert_pptx(src, dst)
        content = dst.read_text()
        assert "dropped unsupported image: unreadable" in content

    def test_unreadable_appears_in_summary(
        self,
        tmp_path: Path,
        monkeypatch: pytest.MonkeyPatch,
        caplog: pytest.LogCaptureFixture,
    ) -> None:
        import logging

        from mdd.convert.pptx import convert_pptx

        src = _make_pptx(tmp_path, [{"title": "Bad pic", "pictures": [_minimal_png()]}])
        dst = tmp_path / "test.pptx.md"
        self._patch_picture_image_raises(monkeypatch, ValueError("no embedded image"))

        mdd_log = logging.getLogger("mdd.convert.pptx")
        mdd_log.addHandler(caplog.handler)
        try:
            with caplog.at_level(logging.WARNING, logger="mdd.convert.pptx"):
                convert_pptx(src, dst)
        finally:
            mdd_log.removeHandler(caplog.handler)
        msgs = " ".join(r.getMessage() for r in caplog.records)
        assert "dropped 1 embedded image" in msgs
        assert "unreadable" in msgs


class TestSkippedShapeWarnings:
    """#92: roll up per-shape unsupported-type warnings into a single summary."""

    def _make_with_line(self, tmp_path: Path) -> Path:
        """Build a pptx whose shape loop emits LINE shapes that aren't handled."""
        import pptx  # pyright: ignore[reportMissingModuleSource]
        from pptx.enum.shapes import MSO_CONNECTOR  # pyright: ignore[reportMissingModuleSource]
        from pptx.util import Inches  # pyright: ignore[reportMissingModuleSource]

        prs: Any = pptx.Presentation()  # pyright: ignore[reportAny]
        slide: Any = prs.slides.add_slide(prs.slide_layouts[6])
        # Two connectors (LINE shapes) so the rollup gets count > 1
        line = MSO_CONNECTOR.STRAIGHT
        slide.shapes.add_connector(line, Inches(1), Inches(1), Inches(2), Inches(2))
        slide.shapes.add_connector(line, Inches(1), Inches(2), Inches(2), Inches(3))
        path = tmp_path / "lines.pptx"
        prs.save(str(path))
        return path

    def test_no_per_shape_warnings_by_default(
        self, tmp_path: Path, capsys: pytest.CaptureFixture[str]
    ) -> None:
        from mdd.convert.pptx import convert_pptx

        src = self._make_with_line(tmp_path)
        dst = tmp_path / "lines.pptx.md"
        convert_pptx(src, dst)
        err = capsys.readouterr().err
        # Old behaviour: one "Warning: slide N: unsupported shape type" line per shape.
        assert "unsupported shape type" not in err

    def test_summary_includes_rollup_by_type(
        self, tmp_path: Path, caplog: pytest.LogCaptureFixture
    ) -> None:
        import logging

        from mdd.convert.pptx import convert_pptx

        src = self._make_with_line(tmp_path)
        dst = tmp_path / "lines.pptx.md"
        mdd_log = logging.getLogger("mdd.convert.pptx")
        mdd_log.addHandler(caplog.handler)
        try:
            with caplog.at_level(logging.WARNING, logger="mdd.convert.pptx"):
                convert_pptx(src, dst)
        finally:
            mdd_log.removeHandler(caplog.handler)
        msgs = " ".join(r.getMessage() for r in caplog.records)
        # New behaviour: a single summary mentioning the type and a count.
        assert "shape(s) skipped" in msgs
        assert "LINE" in msgs
        assert "×2" in msgs


class TestConvertPptxCorruptSource:
    """Issue #129: convert_pptx raises CorruptSourceError for empty/garbage input."""

    def test_zero_byte_pptx_raises_corrupt(self, tmp_path: Path) -> None:
        from mdd.convert import CorruptSourceError
        from mdd.convert.pptx import convert_pptx

        src = tmp_path / "empty.pptx"
        src.write_bytes(b"")  # 0 bytes
        dst = tmp_path / "empty.pptx.md"
        with pytest.raises(CorruptSourceError, match="empty file"):
            convert_pptx(src, dst)
        assert not dst.exists()

    def test_not_a_zip_pptx_raises_corrupt(self, tmp_path: Path) -> None:
        from mdd.convert import CorruptSourceError
        from mdd.convert.pptx import convert_pptx

        src = tmp_path / "tiny.pptx"
        src.write_bytes(b"not a zip, just some random bytes")
        dst = tmp_path / "tiny.pptx.md"
        with pytest.raises(CorruptSourceError, match="not a valid pptx package"):
            convert_pptx(src, dst)
        assert not dst.exists()
