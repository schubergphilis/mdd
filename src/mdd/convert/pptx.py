"""pptx.py — convert .pptx to .pptx.md with one slide per ## section.

Uses python-pptx for shape iteration, image extraction, speaker notes,
and table rendering. Docling is not used for pptx body since python-pptx
provides direct access to all the slide elements (speaker notes, picture
blobs, shape type) that Docling does not surface.
"""

import html
from collections import Counter
from typing import TYPE_CHECKING, Any

import yaml

from mdd.convert import CorruptSourceError
from mdd.utils.logging import get_logger

log = get_logger(__name__)

if TYPE_CHECKING:
    from pathlib import Path


def _same_shape(a: Any, b: Any) -> bool:  # pyright: ignore[reportAny]
    """Return True if *a* and *b* wrap the same pptx shape.

    Tolerates None on either side. python-pptx returns fresh Python wrapper
    objects on each shape lookup (``slide.shapes.title`` and iterating
    ``slide.shapes`` produce different instances even when they reference
    the same underlying XML element), so ``is`` comparisons silently fail.
    Compare by ``shape_id`` instead — unique per shape within a slide.
    """
    if a is None or b is None:
        return False
    try:
        return bool(a.shape_id == b.shape_id)  # pyright: ignore[reportAny]
    except AttributeError:
        return False


def _shape_type_label(shape_type: Any) -> str:  # pyright: ignore[reportAny]
    """Return a short human label for a python-pptx shape type.

    Examples: ``MSO_SHAPE_TYPE.LINE`` -> ``"LINE"``; ``None`` -> ``"None"``;
    anything else -> ``repr(shape_type)``.
    """
    if shape_type is None:
        return "None"
    name: Any = getattr(shape_type, "name", None)  # pyright: ignore[reportAny]
    if isinstance(name, str) and name:
        return name
    return repr(shape_type)


def _render_text_frame(shape: Any, title_shape: Any) -> list[str]:  # pyright: ignore[reportAny]
    """Render a text-frame shape as markdown lines.

    Skips the shape if it IS the title shape (title already emitted as ##).
    Returns empty list for empty text frames or title shapes.
    """
    if title_shape is not None and _same_shape(shape, title_shape):
        return []

    tf: Any = getattr(shape, "text_frame", None)  # pyright: ignore[reportAny]
    if tf is None:
        return []

    lines: list[str] = []
    for para in tf.paragraphs:  # pyright: ignore[reportAny]
        text: str = str(para.text).strip()  # pyright: ignore[reportAny]
        if not text:
            continue
        level: int = int(para.level)  # pyright: ignore[reportAny]
        if level == 0:
            lines.append(text)
        else:
            indent = "  " * (level - 1)
            lines.append(f"{indent}- {text}")

    return lines


def _cell_has_span(tc: Any) -> bool:  # pyright: ignore[reportAny,reportExplicitAny]
    """Return True if the cell's ``<a:tc>`` element declares a grid- or row-span."""
    # lazy: python-pptx pulls lxml; loaded only on pptx conversion
    from pptx.oxml.ns import qn  # pyright: ignore[reportMissingModuleSource]  # noqa: PLC0415

    return bool(
        tc.get(qn("a:gridSpan"))  # pyright: ignore[reportAny]
        or tc.get("gridSpan")  # pyright: ignore[reportAny]
        or tc.get(qn("a:rowSpan"))  # pyright: ignore[reportAny]
        or tc.get("rowSpan")  # pyright: ignore[reportAny]
    )


def _cell_has_multi_para(cell: Any) -> bool:  # pyright: ignore[reportAny,reportExplicitAny]
    """Return True if the cell's text frame holds more than one non-empty paragraph."""
    paras: list[Any] = [  # pyright: ignore[reportAny]
        p
        for p in cell.text_frame.paragraphs
        if str(p.text).strip()  # pyright: ignore[reportAny]
    ]
    return len(paras) > 1


def _is_simple_table(table: Any) -> bool:  # pyright: ignore[reportAny]
    """Return True if all cells contain only plain text (no merged cells, no multi-para)."""
    return not any(
        _cell_has_span(cell._tc) or _cell_has_multi_para(cell)  # pyright: ignore[reportAny]
        for row in table.rows  # pyright: ignore[reportAny]
        for cell in row.cells  # pyright: ignore[reportAny]
    )


def _render_table_markdown(table: Any) -> list[str]:  # pyright: ignore[reportAny]
    """Render a pptx table as a markdown table."""
    rows_data: list[list[str]] = []
    for row in table.rows:  # pyright: ignore[reportAny]
        row_data: list[str] = [  # pyright: ignore[reportAny]
            str(cell.text_frame.text).strip()
            for cell in row.cells  # pyright: ignore[reportAny]
        ]
        rows_data.append(row_data)

    if not rows_data:
        return []

    # Compute column widths
    col_count = len(rows_data[0])
    widths = [max(len(row[c]) for row in rows_data) for c in range(col_count)]
    widths = [max(w, 3) for w in widths]

    lines: list[str] = []
    for i, row in enumerate(rows_data):
        cells = " | ".join(cell.ljust(widths[c]) for c, cell in enumerate(row))
        lines.append(f"| {cells} |")
        if i == 0:
            sep = " | ".join("-" * widths[c] for c in range(col_count))
            lines.append(f"| {sep} |")
    return lines


def _render_table_html(table: Any) -> list[str]:  # pyright: ignore[reportAny]
    """Render a pptx table as a raw HTML table in a fenced block."""
    html_rows: list[str] = ["<table>"]
    for i, row in enumerate(table.rows):  # pyright: ignore[reportAny]
        html_rows.append("  <tr>")
        tag = "th" if i == 0 else "td"
        for cell in row.cells:  # pyright: ignore[reportAny]
            text: str = html.escape(str(cell.text_frame.text).strip(), quote=True)  # pyright: ignore[reportAny]
            html_rows.append(f"    <{tag}>{text}</{tag}>")
        html_rows.append("  </tr>")
    html_rows.append("</table>")
    return ["```{=html}", *html_rows, "```"]


def _render_table(shape: Any) -> list[str]:  # pyright: ignore[reportAny]
    """Render a table shape as markdown or HTML fallback."""
    table: Any = shape.table  # pyright: ignore[reportAny]
    if _is_simple_table(table):
        return _render_table_markdown(table)
    return _render_table_html(table)


def _try_extract_picture(
    shape: Any,  # pyright: ignore[reportAny]
    attachments_dir: Path,
    cache: dict[str, Path],
    dropped_reasons: Counter[str],
) -> str:
    """Write a picture blob via the shared writer; return its markdown link.

    Returns the markdown link, or a placeholder comment for dropped
    images. Two python-pptx failure modes are caught and recorded
    rather than propagated:

    - ``ValueError`` on ``image.ext`` for unrecognised formats (MPO,
      HEIC, ...) — recorded as the format name parsed from the message.
    - ``ValueError("no embedded image")`` and
      ``AttributeError("'Part' object has no attribute 'image'")`` on
      ``shape.image`` itself, when the picture-shape relationship target
      is missing or points at a non-image Part. Recorded as
      ``"unreadable"``.

    Unknown-extension drops inside :func:`write_image` route through
    the same ``dropped_reasons`` counter via ``on_drop``.

    The blob is content-addressed by ``sha1(blob)[:16]``; subsequent
    references to the same blob in the same conversion reuse the
    already-written file (see :mod:`mdd.convert.images`).
    """
    # test-seam: re-imported here so monkeypatch on
    # ``mdd.convert.images.write_image`` reaches this call site.
    from mdd.convert.images import write_image  # noqa: PLC0415

    try:
        # `shape.image` raises ValueError("no embedded image") or
        # AttributeError("'Part' object has no attribute 'image'") when
        # the picture-shape rel target is missing/non-image. `image.ext`
        # raises ValueError on unrecognised formats (MPO, HEIC, ...).
        image: Any = shape.image  # pyright: ignore[reportAny]
        fmt: str = image.ext  # pyright: ignore[reportAny]
    except (ValueError, AttributeError) as exc:
        message = str(exc)
        reason = _format_from_pptx_error(message) if "got '" in message else "unreadable"
        dropped_reasons[reason] += 1
        return f"<!-- dropped unsupported image: {reason} -->"

    def _record(reason: str) -> None:
        dropped_reasons[reason] += 1

    result = write_image(
        attachments_dir,
        image.blob,  # pyright: ignore[reportAny]
        fmt,
        cache=cache,
        on_drop=_record,
    )
    if result is None:
        return f"<!-- dropped unsupported image: {fmt.upper()} -->"
    rel = f"{attachments_dir.name}/{result.rel_path.as_posix()}"
    return f"![]({rel})"


# python-pptx error-message format adapter, version-pinned: the upstream
# `Image.ext` property raises `ValueError` with a message containing
# ``got 'XXX'`` for unrecognised image blobs (MPO, HEIC, ...). We parse
# the format name out of that message in one acknowledged place.
def _format_from_pptx_error(message: str) -> str:
    """Extract the offending format name from python-pptx's error string.

    The exception message reads, e.g.,
    ``unsupported image format, expected one of: dict_keys([...]), got 'MPO'``.
    Returns ``"MPO"`` in that case; ``"unknown"`` if the parse fails.
    """
    marker = "got '"
    idx = message.rfind(marker)
    if idx == -1:
        return "unknown"
    rest = message[idx + len(marker) :]
    end = rest.find("'")
    return rest[:end] if end > 0 else "unknown"


def _notes_text(slide: Any) -> str:  # pyright: ignore[reportAny]
    """Return non-empty speaker notes text, or empty string.

    Catches only AttributeError (no notes_slide attribute) and KeyError
    (missing XML element). Other exceptions propagate so broken notes XML
    is visible rather than silently dropped.
    """
    try:
        notes_slide: Any = slide.notes_slide  # pyright: ignore[reportAny]
        tf: Any = notes_slide.notes_text_frame  # pyright: ignore[reportAny]
        text: str = str(tf.text).strip()  # pyright: ignore[reportAny]
        return text
    except AttributeError, KeyError:
        return ""


def _render_shape(
    shape: Any,  # pyright: ignore[reportAny]
    title_shape: Any,  # pyright: ignore[reportAny]
    attachments_dir: Path,
    image_cache: dict[str, Path],
    skipped_types: Counter[str],
    dropped_image_reasons: Counter[str],
) -> list[str]:
    """Render a single non-title shape; return its markdown lines."""
    # lazy: python-pptx pulls lxml; loaded only on pptx conversion
    from pptx.enum.shapes import (  # noqa: PLC0415
        MSO_SHAPE_TYPE,  # pyright: ignore[reportMissingModuleSource]
    )

    shape_type: Any = shape.shape_type  # pyright: ignore[reportAny]

    if shape_type == MSO_SHAPE_TYPE.PICTURE:  # pyright: ignore[reportAny]
        link = _try_extract_picture(shape, attachments_dir, image_cache, dropped_image_reasons)
        return [link, ""]

    if getattr(shape, "has_text_frame", False):  # pyright: ignore[reportAny]
        rendered = _render_text_frame(shape, title_shape)
        return [*rendered, ""] if rendered else []

    if getattr(shape, "has_table", False):  # pyright: ignore[reportAny]
        rendered_table = _render_table(shape)
        return [*rendered_table, ""] if rendered_table else []

    skipped_types[_shape_type_label(shape_type)] += 1
    return []


def _render_slide(
    slide_num: int,
    slide: Any,  # pyright: ignore[reportAny]
    attachments_dir: Path,
    image_cache: dict[str, Path],
    skipped_types: Counter[str],
    dropped_image_reasons: Counter[str],
) -> list[str]:
    """Render one slide's markdown lines (including title heading and notes)."""
    title_shape: Any = slide.shapes.title  # pyright: ignore[reportAny]
    slide_title: str = (
        str(title_shape.text).strip()  # pyright: ignore[reportAny]
        if title_shape is not None
        else ""
    )

    heading = slide_title or f"Slide {slide_num}"
    slide_lines: list[str] = [f"## {heading}", ""]

    for shape in slide.shapes:  # pyright: ignore[reportAny]
        # Skip the title shape (already rendered as ##).
        # python-pptx returns fresh wrapper objects on each access, so we
        # compare by shape_id rather than identity (see _same_shape).
        if _same_shape(shape, title_shape):
            continue
        slide_lines.extend(
            _render_shape(
                shape,
                title_shape,
                attachments_dir,
                image_cache,
                skipped_types,
                dropped_image_reasons,
            )
        )

    notes = _notes_text(slide)
    if notes:
        slide_lines.extend(["::: notes", notes, ":::", ""])

    return slide_lines


def _emit_summaries(
    src: Path,
    skipped_types: Counter[str],
    dropped_image_reasons: Counter[str],
) -> None:
    """Emit end-of-deck stderr summaries (rolled-up shapes; dropped images)."""
    if skipped_types:
        total = sum(skipped_types.values())
        breakdown = ", ".join(
            f"{label}×{count}"
            for label, count in sorted(skipped_types.items(), key=lambda kv: (-kv[1], kv[0]))
        )
        log.warning("%d shape(s) skipped (unsupported type): %s", total, breakdown)

    if dropped_image_reasons:
        total_dropped = sum(dropped_image_reasons.values())
        formats = ", ".join(
            f"{fmt}×{count}"
            for fmt, count in sorted(dropped_image_reasons.items(), key=lambda kv: (-kv[1], kv[0]))
        )
        log.warning(
            "%s: dropped %d embedded image(s) that could not be loaded (unsupported format: %s).",
            src.name,
            total_dropped,
            formats,
        )


def convert_pptx(src: Path, dst: Path) -> None:
    """Convert src .pptx to dst .pptx.md with one slide per ## section.

    Raises :class:`mdd.convert.CorruptSourceError` when *src* is empty
    or is not a valid Office Open XML package (e.g. truncated downloads,
    deliberately-broken decks). Callers at the sync-dispatch boundary
    route this to a soft-skip counter rather than the hard-error counter
    (issue #129).
    """
    # lazy: python-pptx pulls lxml; loaded only on pptx conversion
    from pptx import Presentation  # pyright: ignore[reportMissingModuleSource]  # noqa: PLC0415
    from pptx.exc import (  # noqa: PLC0415
        PackageNotFoundError,  # pyright: ignore[reportMissingModuleSource]
    )

    # Cheap pre-check: an empty file cannot be a ZIP container; skip the
    # python-pptx open entirely so the error message is unambiguous.
    if src.stat().st_size == 0:
        raise CorruptSourceError(f"{src}: empty file (0 bytes)")

    try:
        prs: Any = Presentation(str(src))  # pyright: ignore[reportAny]
    except PackageNotFoundError as exc:
        raise CorruptSourceError(f"{src}: not a valid pptx package ({exc})") from exc

    # Build frontmatter from core_properties
    props: Any = prs.core_properties  # pyright: ignore[reportAny]
    title: str = str(getattr(props, "title", "") or "").strip()
    author: str = str(getattr(props, "author", "") or "").strip()
    created: Any = getattr(props, "created", None)  # pyright: ignore[reportAny]
    created_str: str = created.isoformat() if created is not None else ""  # pyright: ignore[reportAny]
    slides: list[Any] = list(prs.slides)  # pyright: ignore[reportAny]
    slide_count = len(slides)

    # Determine attachments directory (sibling of dst, named <dst.stem>-attachments)
    # dst is e.g. /path/to/foo.pptx.md; stem is "foo.pptx"; attachments = foo.pptx-attachments
    attachments_dir = dst.parent / (dst.stem + "-attachments")

    # Build YAML frontmatter using yaml.safe_dump for correct encoding
    fm_data: dict[str, object] = {"slide_count": slide_count}
    if title:
        fm_data["title"] = title
    if author:
        fm_data["author"] = author
    if created_str:
        fm_data["created"] = created_str
    fm_block = yaml.safe_dump({"pptx": fm_data}, default_flow_style=False, allow_unicode=True)
    fm_lines = ["---", *fm_block.rstrip().splitlines(), "---", ""]

    body_parts: list[str] = []
    # Per-conversion blob → relative path cache; reused across slides so
    # the same image referenced from N slides hits one file on disk.
    image_cache: dict[str, Path] = {}
    skipped_types: Counter[str] = Counter()
    dropped_image_reasons: Counter[str] = Counter()

    for slide_num, slide in enumerate(slides, 1):
        body_parts.extend(
            _render_slide(
                slide_num,
                slide,
                attachments_dir,
                image_cache,
                skipped_types,
                dropped_image_reasons,
            )
        )

    _emit_summaries(src, skipped_types, dropped_image_reasons)

    content = "\n".join(fm_lines) + "\n".join(body_parts)

    # Atomic write
    dst.parent.mkdir(parents=True, exist_ok=True)
    tmp = dst.with_suffix(".md.tmp")
    tmp.write_text(content, encoding="utf-8")
    tmp.rename(dst)
